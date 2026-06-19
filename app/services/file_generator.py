"""
file_generator.py — Generate PDF, PPTX, and HTML files from AI chat output.

Returns bytes + a suggested filename for each format.
The router streams these back as FileResponse / Response.
"""

import io
import re
import textwrap
from datetime import datetime, timezone


# ── Helpers ───────────────────────────────────────────────────────────────────

def _safe_filename(title: str, ext: str) -> str:
    slug = re.sub(r"[^\w\s-]", "", title).strip().replace(" ", "_")[:40]
    date = datetime.now(timezone.utc).strftime("%Y%m%d")
    return f"{slug}_{date}.{ext}"


def _strip_markdown(text: str) -> str:
    """Convert basic markdown to plain text for PDF/PPTX."""
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)   # bold
    text = re.sub(r"\*(.+?)\*", r"\1", text)         # italic
    text = re.sub(r"`(.+?)`", r"\1", text)            # inline code
    text = re.sub(r"#{1,6}\s+", "", text)             # headings → plain
    text = re.sub(r"^\s*[-*]\s+", "• ", text, flags=re.MULTILINE)  # bullets
    return text.strip()


# ── PDF ───────────────────────────────────────────────────────────────────────

def generate_pdf(title: str, content: str, workflow_title: str = "") -> tuple[bytes, str]:
    """
    Generate a PDF from markdown-ish text.
    Returns (pdf_bytes, filename).
    """
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable
    from reportlab.lib.enums import TA_LEFT, TA_CENTER

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        leftMargin=2.5 * cm,
        rightMargin=2.5 * cm,
        topMargin=2.5 * cm,
        bottomMargin=2.5 * cm,
    )

    styles = getSampleStyleSheet()
    brand_color = colors.HexColor("#1a1a2e")
    accent_color = colors.HexColor("#e94560")

    title_style = ParagraphStyle(
        "CustomTitle",
        parent=styles["Title"],
        fontSize=22,
        textColor=brand_color,
        spaceAfter=6,
        alignment=TA_LEFT,
    )
    subtitle_style = ParagraphStyle(
        "Subtitle",
        parent=styles["Normal"],
        fontSize=10,
        textColor=colors.grey,
        spaceAfter=12,
    )
    body_style = ParagraphStyle(
        "Body",
        parent=styles["Normal"],
        fontSize=11,
        leading=16,
        spaceAfter=8,
        textColor=colors.HexColor("#333333"),
    )
    heading_style = ParagraphStyle(
        "Heading",
        parent=styles["Heading2"],
        fontSize=14,
        textColor=brand_color,
        spaceBefore=14,
        spaceAfter=6,
        borderPad=0,
    )
    bullet_style = ParagraphStyle(
        "Bullet",
        parent=body_style,
        leftIndent=16,
        firstLineIndent=-16,
        spaceAfter=4,
    )

    story = []

    # Header
    story.append(Paragraph(title, title_style))
    if workflow_title:
        story.append(Paragraph(f"Workflow: {workflow_title}", subtitle_style))
    generated_at = datetime.now(timezone.utc).strftime("%B %d, %Y at %H:%M UTC")
    story.append(Paragraph(f"Generated: {generated_at}", subtitle_style))
    story.append(HRFlowable(width="100%", thickness=2, color=accent_color, spaceAfter=16))

    # Parse content line by line
    for line in content.split("\n"):
        stripped = line.strip()
        if not stripped:
            story.append(Spacer(1, 6))
            continue

        if stripped.startswith("# "):
            story.append(Paragraph(stripped[2:], heading_style))
        elif stripped.startswith("## "):
            story.append(Paragraph(stripped[3:], heading_style))
        elif stripped.startswith("### "):
            story.append(Paragraph(stripped[4:], heading_style))
        elif stripped.startswith(("- ", "* ", "• ")):
            text = _strip_markdown(stripped[2:])
            story.append(Paragraph(f"• {text}", bullet_style))
        elif re.match(r"^\d+\. ", stripped):
            text = re.sub(r"^\d+\. ", "", stripped)
            text = _strip_markdown(text)
            story.append(Paragraph(text, bullet_style))
        elif stripped.startswith("**") and stripped.endswith("**"):
            # Bold-only line → treat as a sub-heading
            story.append(Paragraph(_strip_markdown(stripped), heading_style))
        else:
            story.append(Paragraph(_strip_markdown(stripped), body_style))

    doc.build(story)
    return buffer.getvalue(), _safe_filename(title, "pdf")


# ── PPTX ──────────────────────────────────────────────────────────────────────

def generate_pptx(title: str, content: str, workflow_title: str = "") -> tuple[bytes, str]:
    """
    Convert structured text into a PowerPoint presentation.
    Each markdown heading becomes a new slide.
    Returns (pptx_bytes, filename).
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    BRAND = RGBColor(0x1a, 0x1a, 0x2e)
    ACCENT = RGBColor(0xe9, 0x45, 0x60)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GREY = RGBColor(0xF5, 0xF5, 0xF5)

    blank_layout = prs.slide_layouts[6]  # completely blank

    def add_slide(slide_title: str, bullets: list[str], is_title_slide: bool = False):
        slide = prs.slides.add_slide(blank_layout)

        # Background
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = WHITE if not is_title_slide else BRAND

        # Accent bar (left edge)
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.12), prs.slide_height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()

        # Title box
        title_top = Inches(0.5) if not is_title_slide else Inches(2.5)
        txBox = slide.shapes.add_textbox(Inches(0.35), title_top, Inches(12.5), Inches(1.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_title
        p.font.bold = True
        p.font.size = Pt(36) if is_title_slide else Pt(28)
        p.font.color.rgb = WHITE if is_title_slide else BRAND

        if is_title_slide and workflow_title:
            p2 = tf.add_paragraph()
            p2.text = workflow_title
            p2.font.size = Pt(16)
            p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            p3 = tf.add_paragraph()
            p3.text = datetime.now(timezone.utc).strftime("%B %Y")
            p3.font.size = Pt(13)
            p3.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)

        # Content bullets
        if bullets:
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.9), Inches(12.2), Inches(5.2)
            )
            ctf = content_box.text_frame
            ctf.word_wrap = True
            for i, bullet in enumerate(bullets):
                para = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
                clean = _strip_markdown(bullet.lstrip("•-* "))
                para.text = f"• {clean}" if bullet.strip() else ""
                para.font.size = Pt(16)
                para.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                para.space_after = Pt(6)

        return slide

    # Parse content into slides
    # Headings (# / ##) start a new slide; everything else becomes bullets
    slides_data: list[tuple[str, list[str]]] = []
    current_title = title
    current_bullets: list[str] = []

    for line in content.split("\n"):
        stripped = line.strip()
        if not stripped:
            continue
        if re.match(r"^#{1,3}\s+", stripped):
            # Save previous slide
            if current_bullets or current_title != title:
                slides_data.append((current_title, current_bullets))
            current_title = re.sub(r"^#{1,3}\s+", "", stripped)
            current_bullets = []
        else:
            current_bullets.append(stripped)

    # Flush last slide
    if current_title:
        slides_data.append((current_title, current_bullets))

    # Title slide
    add_slide(title, [], is_title_slide=True)

    # Content slides
    for slide_title, bullets in slides_data:
        # Split into max 6 bullets per slide to avoid overflow
        for chunk_start in range(0, max(1, len(bullets)), 6):
            add_slide(slide_title, bullets[chunk_start:chunk_start + 6])

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue(), _safe_filename(title, "pptx")


# ── HTML ──────────────────────────────────────────────────────────────────────

def generate_html(title: str, content: str, workflow_title: str = "") -> tuple[bytes, str]:
    """
    Convert markdown content to a styled standalone HTML file.
    Returns (html_bytes, filename).
    """
    # Very lightweight markdown → HTML conversion (no external deps)
    lines = content.split("\n")
    html_lines = []
    in_ul = False
    in_ol = False

    def close_lists():
        nonlocal in_ul, in_ol
        if in_ul:
            html_lines.append("</ul>")
            in_ul = False
        if in_ol:
            html_lines.append("</ol>")
            in_ol = False

    def inline_md(text: str) -> str:
        text = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", text)
        text = re.sub(r"\*(.+?)\*", r"<em>\1</em>", text)
        text = re.sub(r"`(.+?)`", r"<code>\1</code>", text)
        return text

    for line in lines:
        stripped = line.strip()

        if not stripped:
            close_lists()
            html_lines.append("<br>")
            continue

        if stripped.startswith("# "):
            close_lists()
            html_lines.append(f"<h1>{inline_md(stripped[2:])}</h1>")
        elif stripped.startswith("## "):
            close_lists()
            html_lines.append(f"<h2>{inline_md(stripped[3:])}</h2>")
        elif stripped.startswith("### "):
            close_lists()
            html_lines.append(f"<h3>{inline_md(stripped[4:])}</h3>")
        elif re.match(r"^\d+\. ", stripped):
            if not in_ol:
                close_lists()
                html_lines.append("<ol>")
                in_ol = True
            text = re.sub(r"^\d+\. ", "", stripped)
            html_lines.append(f"<li>{inline_md(text)}</li>")
        elif stripped.startswith(("- ", "* ", "• ")):
            if not in_ul:
                close_lists()
                html_lines.append("<ul>")
                in_ul = True
            html_lines.append(f"<li>{inline_md(stripped[2:])}</li>")
        else:
            close_lists()
            html_lines.append(f"<p>{inline_md(stripped)}</p>")

    close_lists()
    body = "\n".join(html_lines)

    generated_at = datetime.now(timezone.utc).strftime("%B %d, %Y at %H:%M UTC")
    html = textwrap.dedent(f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>{title}</title>
  <style>
    :root {{
      --brand: #1a1a2e;
      --accent: #e94560;
      --text: #333;
      --muted: #777;
      --bg: #fff;
      --surface: #f8f8fc;
      --border: #e5e5e5;
    }}
    * {{ box-sizing: border-box; margin: 0; padding: 0; }}
    body {{
      font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
      color: var(--text);
      background: var(--bg);
      line-height: 1.7;
      max-width: 820px;
      margin: 0 auto;
      padding: 2rem 1.5rem 4rem;
    }}
    header {{
      border-bottom: 3px solid var(--accent);
      padding-bottom: 1.5rem;
      margin-bottom: 2rem;
    }}
    header h1 {{
      font-size: 2rem;
      color: var(--brand);
      margin-bottom: 0.25rem;
    }}
    header .meta {{
      font-size: 0.85rem;
      color: var(--muted);
    }}
    h1, h2, h3 {{ color: var(--brand); margin: 1.8rem 0 0.6rem; }}
    h1 {{ font-size: 1.75rem; }}
    h2 {{ font-size: 1.35rem; border-bottom: 1px solid var(--border); padding-bottom: 0.3rem; }}
    h3 {{ font-size: 1.1rem; }}
    p {{ margin: 0.7rem 0; }}
    ul, ol {{ margin: 0.7rem 0 0.7rem 1.5rem; }}
    li {{ margin: 0.3rem 0; }}
    code {{
      background: var(--surface);
      border: 1px solid var(--border);
      border-radius: 4px;
      padding: 0.1em 0.4em;
      font-size: 0.9em;
      font-family: 'SF Mono', 'Fira Code', monospace;
    }}
    strong {{ color: var(--brand); }}
    footer {{
      margin-top: 3rem;
      padding-top: 1rem;
      border-top: 1px solid var(--border);
      font-size: 0.8rem;
      color: var(--muted);
    }}
    @media print {{
      body {{ max-width: none; padding: 1rem; }}
    }}
  </style>
</head>
<body>
  <header>
    <h1>{title}</h1>
    <div class="meta">
      {"<span>Workflow: " + workflow_title + " &nbsp;|&nbsp; </span>" if workflow_title else ""}
      <span>Generated: {generated_at}</span>
    </div>
  </header>
  <main>
{body}
  </main>
  <footer>Generated by Shift AI · Telfaz11</footer>
</body>
</html>""")

    return html.encode("utf-8"), _safe_filename(title, "html")
